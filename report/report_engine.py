
from pptx import Presentation
from pptx.util import Inches as I, Pt
from pathlib import Path
import json
from PIL import Image, ImageDraw, ImageFont
DEFAULT_TEMPLATE = "assets/Kafaa_Template.pptx"
DEFAULT_BRAND = "assets/kafaa_logo.png"
def _add_title(slide, text):
    try:
        slide.shapes.title.text = text
    except Exception:
        tb = slide.shapes.add_textbox(I(0.8), I(0.6), I(8.0), I(0.8)).text_frame
        tb.clear(); p = tb.paragraphs[0]; r = p.add_run(); r.text = text; r.font.size = Pt(28); r.font.bold=True
def build_pptx(payload_path:str, out:str, brand_mode:str="kafaa"):
    payload = json.loads(Path(payload_path).read_text(encoding="utf-8"))
    tpl = DEFAULT_TEMPLATE if Path(DEFAULT_TEMPLATE).exists() else None
    prs = Presentation(tpl) if tpl else Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _add_title(slide, "Operational Excellence Assessment — Summary")
    body = slide.shapes.placeholders[1].text_frame if len(slide.shapes.placeholders)>1 else slide.shapes.add_textbox(I(0.8),I(1.6),I(8.0),I(4.0)).text_frame
    body.clear()
    def add(b): p = body.add_paragraph(); p.text=b; p.level=0
    client = payload.get("client",{}).get("name_en","Client")
    fin = payload.get("financials",{})
    add(f"Client: {client}")
    add(f"Cost reduction target (SAR): {fin.get('cost_reduction_target_sar','—')}")
    add(f"Frozen cash (SAR): {sum([i.get('value',0) for i in payload.get('muda',{}).get('quantification',{}).get('cash',[])])}")
    add(f"Sales opportunity (SAR/yr): {sum([i.get('value',0) for i in payload.get('muda',{}).get('quantification',{}).get('lost_opportunity',[])])}")
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    _add_title(slide2, "Key Observations (PQCDSM)")
    tf = slide2.shapes.placeholders[1].text_frame if len(slide.shapes.placeholders)>1 else slide2.shapes.add_textbox(I(0.8),I(1.2),I(8.0),I(4.5)).text_frame
    tf.clear()
    for o in (payload.get("observations") or [])[:8]:
        p = tf.add_paragraph(); p.text = f"- {o[:200]}"; p.level=0
    Path(out).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    return out
def build_preview(payload_path:str, out_path:str, brand_mode:str='kafaa'):
    payload = json.loads(Path(payload_path).read_text(encoding='utf-8'))
    W,H = 1200, 675
    im = Image.new('RGB',(W,H),(245,247,250))
    d = ImageDraw.Draw(im)
    try:
        f1 = ImageFont.truetype('DejaVuSans-Bold.ttf', 44)
        f2 = ImageFont.truetype('DejaVuSans.ttf', 28)
    except Exception:
        f1 = f2 = None
    d.text((60,60), "Operational Excellence Assessment", fill=(30,30,30), font=f1)
    client = payload.get('client',{}).get('name_en','Client')
    d.text((60,120), f"Client: {client}", fill=(60,60,60), font=f2)
    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    im.save(out_path, format="PNG")
    return out_path
